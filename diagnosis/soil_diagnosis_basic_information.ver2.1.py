# Step-1 soil_chemical_propertiesフォルダにある測定データ（xlsx）から基本情報を読み込む
# Step-2 基本情報を抽出し、PPTXのfield_propertiesを生成する
# Step-3 圃場測定時画像フォルダにある画像データを読み込む
# Step-4 Step-2のfield_prppertiesに該当画像を追加する
# ver1.2・・基本情報の名称変更と特性深度数値とばらつき数値の追加
# ver2.0・・コメントの自動生成を追加（IDに連動したコメント番号のcsvファイルを読み込みからコメント生成）
# ver2.0・・追加：alldfで欠損値がある列名とID番号をプリント
# ver2.1・・化学性グラフの説明文修正

import datetime
import glob
import os
import pprint
import numpy as np
import pandas as pd
import pptx
from pptx import Presentation
from pptx.util import Cm, Pt

def make_index(alldf):
    prs = Presentation()
    # 1ページ目に「タイトル」スライドのレイアウトを指定
    slide_layout_0 = prs.slide_layouts[0]
    slide_1 = prs.slides.add_slide(slide_layout_0)
    # 2ページ目に「目次」スライドのレイアウトを指定
    slide_layout_1 = prs.slide_layouts[5]
    slide_2 = prs.slides.add_slide(slide_layout_1)
    # 3ページ目以降にID数の空ページとtableを追加、スライドレイアウトを指定
    slide_layout_2 = prs.slide_layouts[6]
    # tableレイアウト設定・・ID別情報ページ
    rows = 25
    cols = 4
    # ID数が26以上ある場合、処理中止する
    isvalid = True
    if 0 < len(alldf['ID']) <= 26:
        for k in range(len(alldf['ID']) * 2):
            m = k + 3
            slide_n = 'slide_' + str(m)
            slide_n = prs.slides.add_slide(slide_layout_2)
            shape_n = slide_n.shapes.add_table(17, cols, Cm(1), Cm(1), Cm(23.5), Cm(15))
            table_in_page = shape_n.table
            if k % 2 == 0:
                table_in_page.columns[0].width = Cm(7)
                table_in_page.columns[1].width = Cm(7)
                table_in_page.columns[2].width = Cm(4.75)
                table_in_page.columns[3].width = Cm(4.75)
                merge_cell_left_top1_a = table_in_page.cell(0, 2)
                merge_cell_right_bottom1_a = table_in_page.cell(0, 3)
                merge_cell_left_top1_a.merge(merge_cell_right_bottom1_a)
                merge_cell_left_top1_a.text = "圃場画像"
                merge_cell_left_top2_a = table_in_page.cell(1, 2)
                merge_cell_right_bottom2_a = table_in_page.cell(8, 3)
                merge_cell_left_top2_a.merge(merge_cell_right_bottom2_a)
                merge_cell_left_top3_a = table_in_page.cell(9, 2)
                merge_cell_right_bottom3_a = table_in_page.cell(16, 3)
                merge_cell_left_top3_a.merge(merge_cell_right_bottom3_a)
            else:
                merge_cell_left_top1_b = table_in_page.cell(0, 0)
                merge_cell_right_bottom1_b = table_in_page.cell(0, 1)
                merge_cell_left_top1_b.merge(merge_cell_right_bottom1_b)
                merge_cell_left_top1_b.text = "土壌化学性グラフ"
                merge_cell_left_top2_b = table_in_page.cell(0, 2)
                merge_cell_right_bottom2_b = table_in_page.cell(0, 3)
                merge_cell_left_top2_b.merge(merge_cell_right_bottom2_b)
                merge_cell_left_top2_b.text = "土壌硬度グラフ"
                merge_cell_left_top3_b = table_in_page.cell(10, 0)
                merge_cell_right_bottom3_b = table_in_page.cell(12, 1)
                merge_cell_left_top3_b.merge(merge_cell_right_bottom3_b)
                merge_cell_left_top3_b.text = "【グラフ共通・棒色】基準値（上下限）\n" \
                                              "赤：基準値上限以上、グレ―：基準値内\n" \
                                              "青：基準値下限以下\n" \
                                              "【グラフ共通・縦点線】\n" \
                                              "黄点線：基準値上限100％、赤点線：基準値上限200％"
                merge_cell_left_top3_b.text_frame.paragraphs[0].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[1].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[2].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[3].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[4].font.size = Pt(12)
                merge_cell_left_top4_b = table_in_page.cell(10, 2)
                merge_cell_right_bottom4_b = table_in_page.cell(12, 3)
                merge_cell_left_top4_b.merge(merge_cell_right_bottom4_b)
                merge_cell_left_top4_b.text = "【グラフ共通：水平線】赤破線：作土深の中央値\n" \
                                              "【左グラフ：垂直線】グレー破線：1.5メガパスカル\n" \
                                              "【中グラフ：水平線】赤点線（上下）内：全体68%範囲\n" \
                                              "【左グラフ・棒色（全体）】作土深の中央値の深さ\n" \
                                              "：青：15㎝以下、黄：30㎝以上、グレー：15㎝以上30㎝以下\n" \
                                              "【中・線色】ばらつき度：緑＞グレー＞オレンジ"
                merge_cell_left_top4_b.text_frame.paragraphs[0].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[1].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[2].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[3].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[4].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[5].font.size = Pt(12)
                merge_cell_left_top5_b = table_in_page.cell(13, 0)
                merge_cell_right_bottom5_b = table_in_page.cell(16, 3)
                merge_cell_left_top5_b.merge(merge_cell_right_bottom5_b)
                # merge_cell_left_top5_b.text = "【コメント】"
    else:
        print("ID数が0または26以上あるため、処理を中止します")
        isvalid = False

    # テキストの編集・・「タイトル」スライド
    # 報告日の取得
    d_today = datetime.date.today()
    title = slide_1.placeholders[0]
    title.text = "土壌診断結果報告書"
    subtitle = slide_1.placeholders[1]
    subtitle_text = "Agsoil株式会社\n報告日：" + str(d_today)
    subtitle.text = subtitle_text
    # テキストの編集・・「目次」スライト
    title2 = slide_2.placeholders[0]
    title2.text = "目次"
    # 目次表の生成
    # 表のレイアウト設定・・目次
    rows = 14
    cols = 6
    table_shape = slide_2.shapes.add_table(rows, cols, Cm(1), Cm(3.5), Cm(23), Cm(10))
    table = table_shape.table
    table.columns[0].width = Cm(1.5)
    table.columns[1].width = Cm(5)
    table.columns[2].width = Cm(5)
    table.columns[3].width = Cm(1.5)
    table.columns[4].width = Cm(5)
    table.columns[5].width = Cm(5)

    # 列見出しのテキスト設定
    category = ['No', 'ID', '圃場名']
    alldf_indexs = alldf.loc[:, ['ID', '圃場名']]
    indexs_count = len(alldf_indexs)
    for i in range(len(category)):
        if indexs_count <= 12:
            cell0 = table.cell(0, i)  # cellオブジェクトの取得
            cell0.text = category[i]  # textプロパティで値を設定する
        else:
            cell0 = table.cell(0, i)  # cellオブジェクトの取得
            cell0.text = category[i]  # textプロパティで値を設定する
            cell1 = table.cell(0, i + 3)
            cell1.text = category[i]
    # alldfから目次を作成する、13以上ある場合列を変更する処理追加
    for j, alldf_index in alldf_indexs.iterrows():
        if j <= 12:
            cell0 = table.cell(j + 1, 0)  # cellオブジェクトの取得
            cell0.text = str(j + 1)  # textプロパティで値を設定する
            cell1 = table.cell(j + 1, 1)
            cell1.text = str(alldf_index['ID'])
            pg1 = cell1.text_frame.paragraphs[0]
            pg1.font.size = Pt(16)
            cell2 = table.cell(j + 1, 2)
            cell2.text = alldf_index['圃場名']
            pg2 = cell2.text_frame.paragraphs[0]
            pg2.font.size = Pt(16)
        else:
            cell0 = table.cell(j - 12, 3)  # cellオブジェクトの取得
            cell0.text = str(j + 1)  # textプロパティで値を設定する
            cell1 = table.cell(j - 12, 4)
            cell1.text = str(alldf_index['ID'])
            pg1 = cell1.text_frame.paragraphs[0]
            pg1.font.size = Pt(16)
            cell2 = table.cell(j - 12, 5)
            cell2.text = alldf_index['圃場名']
            pg2 = cell2.text_frame.paragraphs[0]
            pg2.font.size = Pt(16)
    # PowerPointを保存
    prs.save("output/create_powerpnt.pptx")


def make_picture_table(alldf):
    # 土壌物理性グラフ（jpeg）のリストを読み込む
    filedir1 = 'C:/Users/minam/Desktop/soil_physical_graph2/'
    files_pg = glob.glob(filedir1 + '/**/*.jpeg', recursive=True)
    # 土壌化学性グラフ（jpeg）のリストを読み込む
    filedir2 = 'C:/Users/minam/Desktop/soil_chemical_graph2/'
    files_cg = glob.glob(filedir2 + '/**/*.jpeg', recursive=True)
    # 圃場画像（jpeg）のリストを読み込む
    filedir3 = 'C:/Users/minam/Desktop/hojyo_picture2/'
    files_hp = glob.glob(filedir3 + '/**/*.jpeg', recursive=True)
    # alldf.to_csv('ccc.csv', encoding='Shift-jis')

    # ベースDataframe（alldf2）を作成・・ここにグラフ・画像を紐付ける
    id_list = []
    col_list1 = ['ID']
    for id in alldf['ID']:
        id_list.append(id)
    alldf2 = pd.DataFrame(data=id_list, columns=col_list1)
    # alldf2.to_csv('alldf2-0.csv', encoding='Shift-jis')

    # 土壌物理性診断グラフとIDを結合したDataframe（Pg_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    pg_id = []
    pg_list = []
    col_list2 = ['土壌物理性診断グラフ']
    for file_pg in files_pg:
        pg_name = file_pg.split('\\')
        pg_name = pg_name[-1]
        pg_name_id = pg_name.split('_')
        pg_name_id = pg_name_id[1]
        pg_id.append(pg_name_id)
        pg_list.append(file_pg)
    # print(pg_id)
    # print(pg_list)
    pg_id = pd.DataFrame(data=pg_id, columns=col_list1)
    pg_list = pd.DataFrame(data=pg_list, columns=col_list2)
    pg_df = pd.concat([pg_id, pg_list], axis=1)
    alldf2 = pd.merge(alldf2, pg_df, left_on='ID', right_on='ID')
    # alldf2.to_csv('alldf2-1.csv', encoding='Shift-jis')

    # 土壌化学性診断グラフとIDを結合したDataframe（cg_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    cg_id = []
    cg_list = []
    col_list3 = ['土壌化学性診断グラフ']
    for file_cg in files_cg:
        cg_name = file_cg.split('\\')
        cg_name = cg_name[-1]
        cg_name_id = cg_name.split('_')
        cg_name_id = cg_name_id[1]
        cg_id.append(cg_name_id)
        cg_list.append(file_cg)
    cg_id = pd.DataFrame(data=cg_id, columns=col_list1)
    cg_list = pd.DataFrame(data=cg_list, columns=col_list3)
    cg_df = pd.concat([cg_id, cg_list], axis=1)
    alldf2 = pd.merge(alldf2, cg_df, left_on='ID', right_on='ID')
    # alldf2.to_csv('alldf2-2.csv', encoding='Shift-jis')

    # 圃場画像とIDを結合したDataframe（hp_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    # ベースのDataframe（hp_df）を作成
    col_list4 = ['all', 'left', 'right', 'center', 'other']
    hp_df = pd.DataFrame(data=id_list, columns=col_list1)
    # Dataframe（hp_df）に空列を追加
    for col_name in col_list4:
        hp_df[col_name] = np.nan

    # ベースのDataframe（hp_df）に該当する画像パスを代入する
    file_hp_list = {}
    for file_hp in files_hp:
        hp_name = file_hp.split('\\')
        hp_name = hp_name[-1]
        hp_name_parts = hp_name.split('_')
        hp_name_id = hp_name_parts[1]
        hp_position = hp_name_parts[3]
        file_hp_list[str(hp_name_id) + '_' + hp_position] = file_hp
    for i, hp_id in enumerate(hp_df['ID']):
        for col in col_list4:
            if str(hp_id) + '_' + col in file_hp_list:
                hp_df.loc[i, col] = file_hp_list[str(hp_id) + '_' + col]
    alldf2 = pd.merge(alldf2, hp_df, left_on='ID', right_on='ID')
    # alldf2.to_csv('alldf2-3.csv', encoding='Shift-jis')

    # alldfとalldf2をmergeしてalldf_setを生成
    alldfset = pd.merge(alldf, alldf2, left_on='ID', right_on='ID')
    # alldfset.to_csv('alldfset.csv', encoding='Shift-JIS')
    # alldfsetから基本情報、グラフ・画像情報を抽出しPPTXに代入する関数
    set_basic_information(alldfset)


def get_layout_type(row):
    """
    # Todo: レイアウトに応じて拡張する
    レイアウトのタイプと画像の配列を返す
    :param row:
    """
    layout_type = None
    layout_pictures = []
    if not pd.isna(row['all']):
        print("True", "aaa")
        layout_type = "all"
        layout_pictures = [row['土壌化学性診断グラフ'], row['土壌物理性診断グラフ'], row['all']]
    elif not pd.isna(row['left']) and not pd.isna(row['right']):
        print("True", "bbb")
        layout_type = "left_and_right"
        layout_pictures = [row['土壌化学性診断グラフ'], row['土壌物理性診断グラフ'],
                           row['left'], row['right']]
    else:
        raise ValueError("対象の圃場画像が見つかりません")

    return layout_type, layout_pictures


def set_basic_information(alldfset):
    prs = pptx.Presentation("output/create_powerpnt.pptx")
    # 基本情報で代入する項目のみを抽出したdataframe（alldataset1）を生成
    alldfset1 = alldfset.iloc[:, 1:16]
    # alldfset1.to_csv("alldataset1.csv", encoding='Shift-JIS')
    # alldfset.to_csv('bbb.csv', encoding='Shift-jis')

    headers = 2
    id_pages = 2  # 1圃場当たりに必要なページ数
    for i in range(len(alldfset1)):
        page = i * id_pages + headers
        table_in_page = prs.slides[page].shapes[0].table
        alldfset1_row = alldfset1.iloc[i, :]
        alldfset1_col = alldfset1_row.index
        for k, col_name in enumerate(alldfset1_col):
            col_value = alldfset1_row[col_name]
            table_in_page.cell(k, 0).text = str(col_name)
            table_in_page.cell(k, 1).text = str(col_value)
            pg2 = table_in_page.cell(k, 1).text_frame.paragraphs[0]
            pg2.font.size = Pt(16)
    # 圃場画像の貼り付け
    # 圃場画像URLを抽出したdataframe（alldfset2）を生成
    alldfset2 = alldfset.iloc[:, 22:]
    # alldfset2.to_csv("alldataset2.csv", encoding='Shift-JIS')
    # 圃場画像の情報を読み込み

    for index, row in alldfset2.iterrows():
        layout_type, layout_pictures = get_layout_type(row)
        page_1 = index * 2 + 2
        page_2 = index * 2 + 3
        slide_in_page_1 = prs.slides[page_1]
        slide_in_page_2 = prs.slides[page_2]
        table_in_page_1 = prs.slides[page_1].shapes[0].table
        table_in_page_2 = prs.slides[page_2].shapes[0].table
        if layout_type == "all":
            pic_top_1 = Cm(3.25)
            pic_left_1 = Cm(15)
            pic_height_1 = Cm(9)
            slide_in_page_1.shapes.add_picture(layout_pictures[2],
                                               pic_left_1, pic_top_1, pic_height_1)
            picture_cell1 = table_in_page_1.cell(1, 2)
            picture_cell1.text = "圃場全体"
            pic_left_ch = Cm(1)
            pic_left_py = Cm(12.75)
            pic_top_ch = Cm(2)
            pic_height_ch = Cm(11.7)
            slide_in_page_2.shapes.add_picture(layout_pictures[0],
                                               pic_left_ch, pic_top_ch, pic_height_ch)
            slide_in_page_2.shapes.add_picture(layout_pictures[1],
                                               pic_left_py, pic_top_ch, pic_height_ch)
        else:
            pic_top_1 = Cm(3.25)
            pic_top_2 = Cm(11.75)
            pic_left_1 = Cm(15)
            pic_height_1 = Cm(9)
            slide_in_page_1.shapes.add_picture(layout_pictures[2],
                                               pic_left_1, pic_top_1, pic_height_1)
            slide_in_page_1.shapes.add_picture(layout_pictures[3],
                                               pic_left_1, pic_top_2, pic_height_1)
            picture_cell1 = table_in_page_1.cell(1, 2)
            picture_cell2 = table_in_page_1.cell(9, 2)
            picture_cell1.text = "圃場左側から中心"
            picture_cell2.text = "圃場中心から右側"
            pic_left_ch = Cm(1)
            pic_left_py = Cm(12.75)
            pic_top_ch = Cm(2)
            pic_height_ch = Cm(11.7)
            slide_in_page_2.shapes.add_picture(layout_pictures[0],
                                               pic_left_ch, pic_top_ch, pic_height_ch)
            slide_in_page_2.shapes.add_picture(layout_pictures[1],
                                               pic_left_py, pic_top_ch, pic_height_ch)
    # コメントに代入する項目のみを抽出したdataframe（alldataset3）を生成
    alldfset3 = alldfset.iloc[:, 16:22]
    # alldfset3.to_csv("alldataset3.csv", encoding='Shift-JIS')
    for index, rows in alldfset3.iterrows():
        rows.dropna(inplace=True)
        page_no = index * 2 + 3
        table_in_page = prs.slides[page_no].shapes[0].table
        comment_cell = table_in_page.cell(13, 0)
        textframe = comment_cell.text_frame
        comment = '\n'.join(rows)
        textframe.paragraphs[0].text = "【コメント】" + comment
        textframe.paragraphs[0].font.size = Pt(15)
    # PowerPointを保存
    # 報告日の取得
    d_today = datetime.date.today()
    isvalid = True
    # alldfset.to_csv('aaa.csv', encoding='Shift-jis')
    if len(set(alldfset['出荷団体名'])) == 1:
        group_name = str(alldfset['出荷団体名'][0])
        save_prs_dir = "C:/Users/minam/Desktop/soil_analysisi_report_save/"
        save_prs_name = save_prs_dir + "土壌診断報告書_" + group_name + '_' + str(d_today) + '.pptx'
        print(save_prs_name)
        prs.save(save_prs_name)
    else:
        print("出荷団体名が複数あります。PPTXの保存を中断しました。")
        isvalid = False


def exchange_comment(df_comment):
    comment_dict = {'PH_1': "PHが7以上（アルカリ性）で、PHを矯正する必要があります。",
                    'PH_2': "PHが6.5以上（弱アルカリ性）で、微量要素欠乏を起こす可能性があります。",
                    'PH_3': "PHが6以下（酸性）で、PHを矯正する必要があります。",
                    'PH_4': "PHが低い（酸性）のは窒素過多が原因なので、窒素の施用を控えてください。",
                    'PH_5': "PHは適正範囲（6.0-6.5）です。",
                    'N_1': "窒素は十分残っており、減肥可能です。",
                    'N_2': "ECが高いので、窒素の減肥が可能です。",
                    'N_3': "窒素は過剰で、これ以上の施用は不要です。",
                    'N_4': "窒素および塩基類は過剰で、これ以上の施用は不要です。",
                    'N_5': "EC・無機態窒素の残量は基準値上限以内です。",
                    'P_1': "リン酸は十分残っており、減肥可能です。",
                    'P_2': "リン酸は過剰で、これ以上の施用は不要です。",
                    'P_3': "リン酸吸収力が低いため、施肥量は通常の1.5倍程度必要です。",
                    'P_4': "リン酸吸収力が低いため、施肥量は通常の2.0倍程度必要です。",
                    'P_5': "リン酸吸収力が低いため、施肥量は通常の3.0倍程度必要です。",
                    'P_6': "PHが低い（酸性）ため、リン酸が効きにくくなっています。",
                    'P_7': "リン酸の残量は基準値上限以内です。",
                    'ENKI_1': "塩基飽和度が高いので、塩基類を施用しても土が保持できません",
                    'ENKI_2': "塩基飽和度が高いので、塩基類を施用しても土が保持できません。",
                    'ENKI_3': "カリは過剰で、これ以上の施用は不要です。",
                    'ENKI_4': "塩基バランス（苦土不足、カリ過剰）が崩れています。",
                    'ENKI_5': "塩基バランス（カルシウム不足、苦土・カリ過剰）が崩れています。",
                    'ENKI_7': "塩基バランス（苦土過剰）が崩れています。",
                    'SP_1': "保肥力が低いため、追肥型の施肥を心がけてください。",
                    'SP_2': "十分な腐食があり、高温時の栽培では窒素の減肥が可能です。",
                    'SP_3': "腐食は3%以上（8%以内）あります。",
                    'koudo_1': "作土深は浅め（20ｃｍ以内）で、圃場内での作土深は比較的揃っています。",
                    'koudo_2': "作土深は浅め（20ｃｍ以内）で、圃場内での作土深はややばらついています。",
                    'koudo_3': "作土深は浅め（20ｃｍ以内）で、圃場内での作土深のばらつきは普通です。",
                    'koudo_4': "作土深は深め（30ｃｍ以状）で、圃場内での作土深は比較的揃っています。",
                    'koudo_5': "作土深は深め（30ｃｍ以上）で、圃場内での作土深はややばらついています。",
                    'koudo_6': "作土深は深め（30ｃｍ以上）で、圃場内での作土深のばらつきは普通です。",
                    'koudo_7': "作土深は普通（20-30ｃｍ）で、圃場内での作土深は比較的揃っています。",
                    'koudo_8': "作土深は普通（20-30ｃｍ）で、圃場内での作土深はややばらついています。",
                    'koudo_9': "作土深は普通（20-30ｃｍ）で、圃場内での作土深のばらつきは普通です。"
                    }
    df_comment.replace(comment_dict, inplace=True)
    # df_comment.to_csv("df_comment.csv", encoding='Shift-JIS')


if __name__ == '__main__':
    # Step-1 フォルダにある測定データ（xlsx）を読み込む
    # 読み込むデータを特定するfolder-pathをfiledirに設定する
    filedir = 'C:/Users/minam/Desktop/soil_chemical_properties/'
    # フォルダー内にあるフォルダー名をfolderlist、ファイル名をfilesに所得する
    folderlist = os.listdir(filedir)
    files = glob.glob(filedir + '/**/*.xlsx', recursive=True)
    # コメントCSVファイルの読み込み
    files2 = glob.glob(filedir + '/**/*.csv', recursive=True)
    pprint.pprint(files2)

    # 【Step-1-1】フォルダにある測定データ（.xlsx）から基本情報を読み込む
    for file, file2 in zip(files, files2):
        df_comment = pd.read_csv(file2, index_col=0)
        df_comment = df_comment.sort_values(by="ID")
        # コメント番号からコメント文書に置換する関数
        exchange_comment(df_comment)
        df = pd.read_excel(file, sheet_name='基本情報')
        df = df.loc[:, ['ID', '出荷団体名', '生産者名', '圃場名', '面積（平方メートル）', '圃場位置(緯度)', '圃場位置(経度)', '品目名', '作型']]
        # 「土壌化学性データ」シートから必要情報の取得
        df2 = pd.read_excel(file, sheet_name='土壌化学性データ')
        df2 = df2.loc[:, ['ID', '採土日', '採土法', '仮比重']]
        # 「土壌物理性データ」シートから必要情報の取得
        df3 = pd.read_excel(file, sheet_name='土壌物理性データ')
        df3 = df3.loc[:, ['ID', '測定日', '測定法', '測定状態', '作土深の中心値']]
        # 【Step-1-2】取得データの結合（キー列'ID'）と欠損値の判定
        alldf = pd.merge(pd.merge(df, df2, left_on='ID', right_on='ID'), df3,
                         left_on='ID', right_on='ID')
        # alldfと文書変換したdf_commentをmergeする
        alldf = pd.merge(alldf, df_comment, left_on='ID', right_on='ID')

        isvalid = True
        for i in range(len(alldf)):
            if alldf.loc[i].isnull().any():
                print(i, "行目", "欠損値のある行が含まれています", "ID:", alldf.loc[i][0])
                isvalid = False
            else:
                print(i, "行目", "必要情報は正常です")
        # IDをキー列にして昇順ソート、indexを振り直す　注：欠損値あっても行削除せず
        # 【課題】sortができていない？
        alldf = alldf.sort_values(by="ID")
        # ひな形のPPTXを読み込み目次を作成
        make_index(alldf)
        # ID別の紐付け情報（グラフ2種、圃場画像１～４種）のDataframeを作成する
        make_picture_table(alldf)
