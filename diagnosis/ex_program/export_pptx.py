from pptx import Presentation
from pptx.util import Cm
from diagnosis.domain.field_property import FieldProperty

LAYOUT_BLANK = 6


if __name__ == '__main__':
    # https://www.shibutan-bloomers.com/python-libraly-pptx-5/1188/
    # 前提：output にあるjpegのペアを export_pptx_ordering に書いてから実行
    # TODO: textboxの外枠の線

    with open("export_pptx_ordering.txt", "r", encoding='utf8') as tf:
        orders = tf.read().split('\n')

    pptx = Presentation()

    for order in orders:
        # 3パラメータのときのみ実行される
        if len(order.split(',')) == 3:
            # パラメータ分解
            page, left_picture, right_picture = order.split(',')

            slide = pptx.slides.add_slide(pptx.slide_layouts[LAYOUT_BLANK])

            # 圃場情報
            field_property = FieldProperty()
            field_property.shipping_organization = '県西つくば協同組合'
            field_property.company = '(有)アグリファクトリー'
            field_property.name = '飯田裏'
            field_property.clops = 'ハクサイ'
            field_property.class_clops = '露地'
            field_property.mining_soil_day = '2022/06/13'
            field_property.mining_soil_method = '5点法'
            field_property.measurement_date = '2022/06/13'
            field_property.measurement_method = '5分割×5点=25ポイント'
            text_box = slide.shapes.add_textbox(0, 0, 400, 200)
            text_box.text_frame.text = field_property.formatted_text()

            # 写真右上
            slide.shapes.add_picture(f"data/company_sample.png", Cm(17), Cm(1), Cm(6), Cm(6))
            # 写真左 x, y, w, h
            slide.shapes.add_picture(f"output/{left_picture}", Cm(0.5), Cm(4), Cm(12), Cm(10))
            # 写真右 x, y, w, h
            slide.shapes.add_picture(f"output/{right_picture}", Cm(13), Cm(4), Cm(12), Cm(10))
            # コメント
            text_box = slide.shapes.add_textbox(Cm(1), Cm(15), Cm(17), Cm(17))
            text_box.text_frame.text = '＜コメント＞'

    save_path = '../output/sample.pptx'

    try:
        pptx.save(save_path)
        print('エクスポートが完了しました')
    except PermissionError:
        print('pptxが開かれていたため、エクスポートは失敗しました')

