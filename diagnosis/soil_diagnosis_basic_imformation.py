# Step-1 soil_chemical_propertiesフォルダにある測定データ（xlsx）から基本情報を読み込む
# Step-2 基本情報を抽出し、PPTXのfield_propertiesを生成する
# Step-3 圃場測定時画像フォルダにある画像データを読み込む
# Step-4 Step-2のfield_prppertiesに該当画像を追加する
import datetime
import glob
import math
import os
import pprint

import numpy as np
import pandas as pd
import pptx
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_AUTO_SIZE


def make_index(alldf):
    prs = Presentation()
    # 1ページ目に「タイトル」スライドのレイアウトを指定
    slide_layout_0 = prs.slide_layouts[0]
    slide_1 = prs.slides.add_slide(slide_layout_0)
    # 2ページ目に「目次」スライドのレイアウトを指定
    slide_layout_1 = prs.slide_layouts[5]
    slide_2 = prs.slides.add_slide(slide_layout_1)
    # 3ページ目以降にID数の空ページとtableを追加、スライドレイアウトを指定
    slide_layout_2 = prs.slide_layouts[6]
    # tableレイアウト設定・・ID別情報ページ
    rows = 25
    cols = 4
    # table_shape2 = slide_2.shapes.add_table(rows, cols, Cm(1), Cm(3), Cm(15), Cm(15))
    # ID数が26以上ある場合、処理中止する
    isvalid = True
    if 0 < len(alldf['ID']) <= 26:
        for k in range(len(alldf['ID']) * 2):
            m = k + 3
            slide_n = 'slide_' + str(m)
            # print(slide_n)
            slide_n = prs.slides.add_slide(slide_layout_2)
            shape_n = slide_n.shapes.add_table(17, cols, Cm(1), Cm(1), Cm(23.5), Cm(15))
            table_in_page = shape_n.table
            if k % 2 == 0:
                table_in_page.columns[0].width = Cm(7)
                table_in_page.columns[1].width = Cm(7)
                table_in_page.columns[2].width = Cm(4.75)
                table_in_page.columns[3].width = Cm(4.75)
                merge_cell_left_top1_a = table_in_page.cell(0, 2)
                merge_cell_right_bottom1_a = table_in_page.cell(0, 3)
                merge_cell_left_top1_a.merge(merge_cell_right_bottom1_a)
                merge_cell_left_top1_a.text = "圃場画像"
                merge_cell_left_top2_a = table_in_page.cell(1, 2)
                merge_cell_right_bottom2_a = table_in_page.cell(8, 3)
                merge_cell_left_top2_a.merge(merge_cell_right_bottom2_a)
                # merge_cell_left_top2_a.text = str(1)
                merge_cell_left_top3_a = table_in_page.cell(9, 2)
                merge_cell_right_bottom3_a = table_in_page.cell(16, 3)
                merge_cell_left_top3_a.merge(merge_cell_right_bottom3_a)
                # merge_cell_left_top3_a.text = str(2)
            else:
                merge_cell_left_top1_b = table_in_page.cell(0, 0)
                merge_cell_right_bottom1_b = table_in_page.cell(0, 1)
                merge_cell_left_top1_b.merge(merge_cell_right_bottom1_b)
                merge_cell_left_top1_b.text = "土壌化学性グラフ"
                merge_cell_left_top2_b = table_in_page.cell(0, 2)
                merge_cell_right_bottom2_b = table_in_page.cell(0, 3)
                merge_cell_left_top2_b.merge(merge_cell_right_bottom2_b)
                merge_cell_left_top2_b.text = "土壌硬度グラフ"
                merge_cell_left_top3_b = table_in_page.cell(10, 0)
                merge_cell_right_bottom3_b = table_in_page.cell(12, 1)
                merge_cell_left_top3_b.merge(merge_cell_right_bottom3_b)
                merge_cell_left_top3_b.text = "【棒色】赤：基準値上限以上、グレ―：基準値内\n" \
                                              "                 青：基準値下限以下\n" \
                                              "【縦点線】\n" \
                                              "黄：基準値上限100%、赤：基準値上限200%"
                merge_cell_left_top3_b.text_frame.paragraphs[0].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[1].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[2].font.size = Pt(12)
                merge_cell_left_top3_b.text_frame.paragraphs[3].font.size = Pt(12)
                merge_cell_left_top4_b = table_in_page.cell(10, 2)
                merge_cell_right_bottom4_b = table_in_page.cell(12, 3)
                merge_cell_left_top4_b.merge(merge_cell_right_bottom4_b)
                merge_cell_left_top4_b.text = "【共通：水平線】赤破線：作土深の中央値\n" \
                                              "【共通：垂直線】黒点線：1.5メガパスカル\n" \
                                              "【中：水平線】赤点線：全体68%範囲\n" \
                                              "【左・棒色】作土深の深さ：青：15㎝以下、黄：30㎝以上" \
                                              "\n【中・線色】ばらつき度：緑＞グレー＞橙"
                merge_cell_left_top4_b.text_frame.paragraphs[0].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[1].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[2].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[3].font.size = Pt(12)
                merge_cell_left_top4_b.text_frame.paragraphs[4].font.size = Pt(12)
                merge_cell_left_top5_b = table_in_page.cell(13, 0)
                merge_cell_right_bottom5_b = table_in_page.cell(16, 3)
                merge_cell_left_top5_b.merge(merge_cell_right_bottom5_b)
                merge_cell_left_top5_b.text = "【コメント】"
    else:
        print("ID数が0または26以上あるため、処理を中止します")
        isvalid = False

    # テキストの編集・・「タイトル」スライド
    # 報告日の取得
    d_today = datetime.date.today()
    title = slide_1.placeholders[0]
    title.text = "土壌診断結果報告書"
    subtitle = slide_1.placeholders[1]
    subtitle_text = "Agsoil株式会社\n報告日：" + str(d_today)
    subtitle.text = subtitle_text
    # テキストの編集・・「目次」スライト
    title2 = slide_2.placeholders[0]
    title2.text = "目次"
    # title2.width = Cm(1)
    # 目次表の生成
    # 表のレイアウト設定・・目次
    # rows = len(alldf) + 1
    rows = 14
    # cols = 3
    cols = 6
    table_shape = slide_2.shapes.add_table(rows, cols, Cm(1), Cm(3.5), Cm(23), Cm(10))
    # table_shape = slide_2.shapes.add_table(rows, cols, Cm(3), Cm(3), Cm(15), Cm(5))
    table = table_shape.table
    table.columns[0].width = Cm(1.5)
    table.columns[1].width = Cm(5)
    table.columns[2].width = Cm(5)
    table.columns[3].width = Cm(1.5)
    table.columns[4].width = Cm(5)
    table.columns[5].width = Cm(5)

    # 列見出しのテキスト設定
    category = ['No', 'ID', '圃場名']
    for i in range(len(category)):
        cell = table.cell(0, i)  # cellオブジェクトの取得
        cell.text = category[i]  # textプロパティで値を設定する
    # alldfから目次を作成する
    alldf_indexs = alldf.loc[:, ['ID', '圃場名']]
    for j, alldf_index in alldf_indexs.iterrows():
        cell0 = table.cell(j + 1, 0)  # cellオブジェクトの取得
        cell0.text = str(j + 1)  # textプロパティで値を設定する
        cell1 = table.cell(j + 1, 1)
        cell1.text = str(alldf_index['ID'])
        cell2 = table.cell(j + 1, 2)
        cell2.text = alldf_index['圃場名']
    # PowerPointを保存
    prs.save("output/create_powerpnt.pptx")


def make_picture_table(alldf):
    # 土壌物理性グラフ（jpeg）のリストを読み込む
    filedir1 = 'C:/Users/minam/Desktop/soil_physical_graph2/'
    files_pg = glob.glob(filedir1 + '/**/*.jpeg', recursive=True)
    # pprint.pprint(files_pg)
    # 土壌化学性グラフ（jpeg）のリストを読み込む
    filedir2 = 'C:/Users/minam/Desktop/soil_chemical_graph2/'
    files_cg = glob.glob(filedir2 + '/**/*.jpeg', recursive=True)
    # pprint.pprint(files_cg)
    # 圃場画像（jpeg）のリストを読み込む
    filedir3 = 'C:/Users/minam/Desktop/hojyo_picture2/'
    files_hp = glob.glob(filedir3 + '/**/*.jpeg', recursive=True)
    # pprint.pprint(files_hp)

    # ベースDataframe（alldf2）を作成・・ここにグラフ・画像を紐付ける
    id_list = []
    col_list1 = ['ID']
    for id in alldf['ID']:
        id_list.append(id)
    alldf2 = pd.DataFrame(data=id_list, columns=col_list1)

    # 土壌物理性診断グラフとIDを結合したDataframe（Pg_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    pg_id = []
    pg_list = []
    col_list2 = ['土壌物理性診断グラフ']
    for file_pg in files_pg:
        pg_name = file_pg.split('\\')
        pg_name = pg_name[-1]
        pg_name_id = pg_name.split('_')
        pg_name_id = pg_name_id[1]
        pg_id.append(pg_name_id)
        pg_list.append(file_pg)
    pg_id = pd.DataFrame(data=pg_id, columns=col_list1)
    pg_list = pd.DataFrame(data=pg_list, columns=col_list2)
    pg_df = pd.concat([pg_id, pg_list], axis=1)
    alldf2 = pd.merge(alldf2, pg_df, left_on='ID', right_on='ID')

    # 土壌化学性診断グラフとIDを結合したDataframe（cg_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    cg_id = []
    cg_list = []
    col_list3 = ['土壌化学性診断グラフ']
    for file_cg in files_cg:
        cg_name = file_cg.split('\\')
        cg_name = cg_name[-1]
        cg_name_id = cg_name.split('_')
        cg_name_id = cg_name_id[1]
        cg_id.append(cg_name_id)
        cg_list.append(file_cg)
    cg_id = pd.DataFrame(data=cg_id, columns=col_list1)
    cg_list = pd.DataFrame(data=cg_list, columns=col_list3)
    cg_df = pd.concat([cg_id, cg_list], axis=1)
    alldf2 = pd.merge(alldf2, cg_df, left_on='ID', right_on='ID')

    # 圃場画像とIDを結合したDataframe（hp_df)を作成し、ベースDataframe（alldf2）に結合（IDキー）
    # ベースのDataframe（hp_df）を作成
    col_list4 = ['all', 'left', 'right', 'center', 'other']
    hp_df = pd.DataFrame(data=id_list, columns=col_list1)
    # Dataframe（hp_df）に空列を追加
    for col_name in col_list4:
        hp_df[col_name] = np.nan

    # ベースのDataframe（hp_df）に該当する画像パスを代入する
    file_hp_list = {}
    for file_hp in files_hp:
        hp_name = file_hp.split('\\')
        hp_name = hp_name[-1]
        hp_name_parts = hp_name.split('_')
        hp_name_id = hp_name_parts[1]
        hp_position = hp_name_parts[3]
        file_hp_list[str(hp_name_id) + '_' + hp_position] = file_hp
    for i, hp_id in enumerate(hp_df['ID']):
        for col in col_list4:
            if str(hp_id) + '_' + col in file_hp_list:
                hp_df.loc[i, col] = file_hp_list[str(hp_id) + '_' + col]
    alldf2 = pd.merge(alldf2, hp_df, left_on='ID', right_on='ID')

    # alldfとalldf2をmergeしてalldf_setを生成
    alldfset = pd.merge(alldf, alldf2, left_on='ID', right_on='ID')
    alldfset.to_csv('alldfset.csv', encoding='Shift-JIS')
    set_basic_information(alldfset)


def get_layout_type(row):
    """
    # Todo: レイアウトに応じて拡張する
    レイアウトのタイプと画像の配列を返す
    :param row:
    """
    layout_type = None
    layout_pictures = []
    if not pd.isna(row['all']):
        print("True", "aaa")
        layout_type = "all"
        layout_pictures = [row['土壌化学性診断グラフ'], row['土壌物理性診断グラフ'], row['all']]
    elif not pd.isna(row['left']) and not pd.isna(row['right']):
        print("True", "bbb")
        layout_type = "left_and_right"
        layout_pictures = [row['土壌化学性診断グラフ'], row['土壌物理性診断グラフ'],
                           row['left'], row['right']]
    else:
        raise ValueError("対象の圃場画像が見つかりません")

    return layout_type, layout_pictures



def set_basic_information(alldfset):
    # print(alldfset)
    prs = pptx.Presentation("output/create_powerpnt.pptx")
    # 基本情報で代入する項目のみを抽出したdataframe（alldataset1）を生成
    alldfset1 = alldfset.iloc[:, 1:15]
    # alldfset1 = alldfset1.drop(alldfset1.columns[[5, 6]], axis=1)
    alldfset1.to_csv("alldataset1.csv", encoding='Shift-JIS')

    headers = 2
    id_pages = 2  # 1圃場当たりに必要なページ数
    for i in range(len(alldfset1)):
        page = i * id_pages + headers
        table_in_page = prs.slides[page].shapes[0].table
        alldfset1_row = alldfset1.iloc[i, :]
        alldfset1_col = alldfset1_row.index
        for k, col_name in enumerate(alldfset1_col):
            col_value = alldfset1_row[col_name]
            table_in_page.cell(k, 0).text = str(col_name)
            table_in_page.cell(k, 1).text = str(col_value)

    # 圃場画像の貼り付け
    # 圃場画像URLを抽出したdataframe（alldfset2）を生成
    # alldfset2 = alldfset.iloc[:, 17:]
    alldfset2 = alldfset.iloc[:, 15:]
    # alldfset2 = alldfset.drop(alldfset.columns[[2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16]], axis=1)
    alldfset2.to_csv("alldataset2.csv", encoding='Shift-JIS')
    # 圃場画像の情報を読み込み

    for index, row in alldfset2.iterrows():
        layout_type, layout_pictures = get_layout_type(row)
        print(layout_type, layout_pictures)
        page_1 = index * 2 + 2
        page_2 = index * 2 + 3
        slide_in_page_1 = prs.slides[page_1]
        slide_in_page_2 = prs.slides[page_2]
        table_in_page_1 = prs.slides[page_1].shapes[0].table
        table_in_page_2 = prs.slides[page_2].shapes[0].table
        if layout_type == "all":
            pic_top_1 = Cm(3)
            pic_left_1 = Cm(15.25)
            pic_height_1 = Cm(9)
            slide_in_page_1.shapes.add_picture(layout_pictures[2],
                                               pic_left_1, pic_top_1, pic_height_1)
            picture_cell1 = table_in_page_1.cell(1, 2)
            picture_cell1.text = "圃場全体"
            pic_left_ch = Cm(1)
            pic_left_py = Cm(12.75)
            pic_top_ch = Cm(2)
            pic_height_ch = Cm(11.7)
            slide_in_page_2.shapes.add_picture(layout_pictures[0],
                                               pic_left_ch, pic_top_ch, pic_height_ch)
            slide_in_page_2.shapes.add_picture(layout_pictures[1],
                                               pic_left_py, pic_top_ch, pic_height_ch)
        else:
            pic_top_1 = Cm(3)
            pic_top_2 = Cm(11.25)
            pic_left_1 = Cm(15.25)
            pic_height_1 = Cm(9)
            slide_in_page_1.shapes.add_picture(layout_pictures[2],
                                               pic_left_1, pic_top_1, pic_height_1)
            slide_in_page_1.shapes.add_picture(layout_pictures[3],
                                               pic_left_1, pic_top_2, pic_height_1)
            picture_cell1 = table_in_page_1.cell(1, 2)
            picture_cell2 = table_in_page_1.cell(9, 2)
            picture_cell1.text = "圃場左側から中心"
            picture_cell2.text = "圃場中心から右側"
            pic_left_ch = Cm(1)
            pic_left_py = Cm(12.75)
            pic_top_ch = Cm(2)
            pic_height_ch = Cm(11.7)
            slide_in_page_2.shapes.add_picture(layout_pictures[0],
                                               pic_left_ch, pic_top_ch, pic_height_ch)
            slide_in_page_2.shapes.add_picture(layout_pictures[1],
                                               pic_left_py, pic_top_ch, pic_height_ch)

    # PowerPointを保存
    prs.save("output/create_powerpnt2.pptx")


if __name__ == '__main__':
    # Step-1 フォルダにある測定データ（xlsx）を読み込む
    # 読み込むデータを特定するfolder-pathをfiledirに設定する
    filedir = 'C:/Users/minam/Desktop/soil_chemical_properties/'
    # フォルダー内にあるフォルダー名をfolderlist、ファイル名をfilesに所得する
    folderlist = os.listdir(filedir)
    files = glob.glob(filedir + '/**/*.xlsx', recursive=True)
    # pprint.pprint(files)

    # 【Step-1-1】フォルダにある測定データ（.xlsx）から基本情報を読み込む
    for file in files:
        df = pd.read_excel(file, sheet_name='基本情報')
        df = df.loc[:,
             ['ID', '出荷団体名', '生産者名', '圃場名', '面積（平方メートル）', '圃場位置(緯度)',
              '圃場位置(経度)', '品目名', '作型']]
        # 「土壌化学性データ」シートから必要情報の取得
        df2 = pd.read_excel(file, sheet_name='土壌化学性データ')
        df2 = df2.loc[:, ['ID', '採土日', '採土法']]
        # 「土壌物理性データ」シートから必要情報の取得
        df3 = pd.read_excel(file, sheet_name='土壌物理性データ')
        df3 = df3.loc[:, ['ID', '測定日', '測定法', '測定状態']]
        # 【Step-1-2】取得データの結合（キー列'ID'）と欠損値の判定
        alldf = pd.merge(pd.merge(df, df2, left_on='ID', right_on='ID'), df3,
                         left_on='ID', right_on='ID')
        # print(alldf, type(alldf))
        # alldf.to_csv('aaa')
        # IDをキー列にして昇順ソート、インデックスもリセット
        alldf = alldf.sort_values(by="ID")
        alldf = alldf.reset_index()
        # sortができていない【課題】
        isvalid = True
        for i in range(len(alldf)):
            if alldf.loc[i].isnull().any():
                print("欠損値のある行が含まれています")
                isvalid = False
            else:
                print("必要情報は正常です")
                # 【Step-1-3】画像タイトルおよび圃場名、採土日、測定日の取得
                df_title = alldf.loc[i]
                picture_titles = df_title[['ID', '出荷団体名', '生産者名', '圃場名', '品目名', '作型']].values
                picture_title = '基本情報_' + '_'.join(picture_titles)
                hojyomei = df_title[['圃場名']].values
                saidobi = df_title[['採土日']].values
                sokuteibi = df_title[['測定日']].values
        # 欠損値（NAN）のある行を削除し、indexを振り直す
        alldf = alldf.dropna(how='any')
        alldf = alldf.reset_index(drop=True)
        # ひな形のPPTXを読み込み目次を作成
        make_index(alldf)
        # ID別の紐付け情報（グラフ2種、圃場画像１～４種）のDataframeを作成する
        make_picture_table(alldf)
